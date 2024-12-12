from pptx.util import Inches
import os
import random
from pptx import Presentation
from pptx.util import Inches, Pt
from PIL import Image
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from centralizedData import plotsFolder, presentationsFolder


MAC = 1
MX = 1

columns = ['orders_per_eff_online', 'eff_online_rs', 'daily_orders',
           # 'overdue_orders_per_total_orders_b_cancel_rate',
           'imperfect_order_rate_bad_rating_rate',
           'priorityChanges',

           # 'DistributionOrdersDiscounts',
           # 'exposure_per_eff_online_b_p1p2',
           'ted_gmv_r_burn_gmv_b2c_gmv_p2c_gmv',
           ]
columnsCountry = [
    'daily_orders_percentages',
    'daily_orders_nominal',
    'eff_online_rs_percentages',
    'eff_online_rs_nominal',
]

TITLE_START = (1.3, 0.1)
ACTIONTITLE_START = (0.5, TITLE_START[1]+0.75)
SUBTITLE_START = (0.5, ACTIONTITLE_START[1]+0.5)
IMG_TOP = SUBTITLE_START[1]+0.8
SIZE = (800, 700)


def getFilesFromFolder(folder):
    files = []
    for root, dirs, file in os.walk(folder):
        for f in file:
            if '.txt' in f:
                continue
            files.append(f)
    return files


def getGoodFile(folder):
    files = getFilesFromFolder(folder)
    pool = ['BB.png', 'TB.png', 'BT.png', 'TT.png']
    for f in files:
        if f not in pool:
            return f
    pool = ['TT.png']
    return random.choice(pool)


def doGreatComment(file: str, position: tuple = None, slide=None):
    FONT_SIZE = 0.14
    column = file.split('/')[-2]
    with open(file, 'r') as file:
        content = file.read()
    inverse = False
    if column in ["b_cancel_rate", 'bad_rating_rate', 'imperfect_orders', 'imperfect_order_rate']:
        inverse = True
    divided = content.split('@')
    if content == 'The change cannot be calculated over the last 3 months':
        return
    if content == '':
        content = 'No comments'
    if len(divided) == 1:
        add_text(slide, content, font_size=FONT_SIZE,
                 bold=False, color=(0, 0, 0), position=position, vertical=False,
                 BG=(242, 242, 242))
        return
    else:
        firstPart = add_text(slide, divided[0], font_size=FONT_SIZE,
                             bold=False, color=(0, 0, 0), position=position, vertical=False,
                             BG=(242, 242, 242))
        if inverse:
            if divided[1] == "a worse tendency":
                divided[1] = "a better tendency"
            else:
                divided[1] = "a worse tendency"
        if divided[1] == "a worse tendency":
            color = (204, 0, 0)
            add_text(slide, divided[1], font_size=FONT_SIZE, new=False, last_p=firstPart,
                     bold=True, color=color, )
        else:
            color = (127, 169, 108)
            add_text(slide, divided[1], font_size=FONT_SIZE, new=False, last_p=firstPart,
                     bold=True, color=color, )
        add_text(slide, divided[2], font_size=FONT_SIZE, new=False, last_p=firstPart,
                 bold=False, color=(0, 0, 0), BG=(242, 242, 242))
    return


def add_text(slide, text_letters, new=True, last_p=None, font_size=0.5, bold=False, color=(0, 0, 0), position=None, vertical=True, BG=None):
    if position is None and new:
        raise ValueError('Position is required')

    if new:
        # Add new textbox
        text_box = slide.shapes.add_textbox(
            Inches(position[0]), Inches(position[1]), Inches(position[2]), Inches(position[3]))
        text_frame = text_box.text_frame
        text_frame.word_wrap = True
        text_frame.auto_size = False
        text_frame.text = text_letters
        if vertical:
            text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        # set font roboto
        font = text_frame.paragraphs[0].runs[0].font
        font.name = 'Roboto'
        # set background color
        if BG is not None:
            fill = text_box.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(BG[0], BG[1], BG[2])

        # Set font properties
        for paragraph in text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.size = Inches(font_size)
                run.font.name = 'Roboto'
                run.font.bold = bold
                run.font.color.rgb = RGBColor(color[0], color[1], color[2])
                run.vertical_anchor = MSO_ANCHOR.MIDDLE
                if BG is not None:
                    fill = text_box.fill
                    fill.solid()
                    fill.fore_color.rgb = RGBColor(BG[0], BG[1], BG[2])

        # Return last paragraph for chaining
        return text_frame.paragraphs[-1]

    elif not new and last_p is not None:
        # Add text to existing paragraph
        run = last_p.add_run()
        run.text = ' ' + text_letters
        run.font.size = Inches(font_size)
        run.font.name = 'Roboto'
        run.font.bold = bold
        run.font.color.rgb = RGBColor(color[0], color[1], color[2])

        # Return updated paragraph for chaining
        return last_p

    else:
        raise ValueError('Invalid arguments')


def createFolder(folder=str):
    folder = f'{presentationsFolder}/{folder}'
    if not os.path.exists(folder):
        os.makedirs(folder)


def getImagesAndNotes(folderInitial: str, pais: str, prioridad: str, country=False):
    graph_image_paths = []
    graph_text_paths = []
    # Adjust this path as per your actual image
    if not country:
        columnsToUse = columns
    else:
        columnsToUse = columnsCountry
    for graph in columnsToUse:
        if not country:
            folder = f'{folderInitial}/{pais}/p{prioridad}/{graph}'
        else:
            folder = f'{folderInitial}/{pais}/All/{graph}'
        path = getGoodFile(folder)
        graph_image_paths.append(folder+'/'+path)
        graph_text_paths.append(folder+'/'+'note.txt')
    return graph_image_paths, graph_text_paths


def calculateSpace(slideSizeX: float, imgInchSize: float, howMuchi: int, size: tuple = SIZE):
    spaceBetweenEquiSpace = round(
        (slideSizeX-(imgInchSize*howMuchi))/(howMuchi+1), 2)
    spaceBetweenBorderSpace = round((slideSizeX-(imgInchSize*howMuchi))/2, 2)
    space = None
    typeOfSpace = None
    if spaceBetweenEquiSpace < 0.5 and spaceBetweenBorderSpace < 0.5:
        imgInchSize = (slideSizeX-(0.3*2))/howMuchi
        space = 0.3
        typeOfSpace = 'border'
    elif spaceBetweenEquiSpace < 0.5:
        space = spaceBetweenBorderSpace
        typeOfSpace = 'border'
    else:
        space = spaceBetweenEquiSpace
        typeOfSpace = 'equi'
    heightOfImg = imgInchSize*(size[1]/size[0])
    widthOfImg = imgInchSize
    return space, typeOfSpace, widthOfImg, heightOfImg


def addGraphAndText(slide, graph_image_path: str, graph_note_path: str, Yposition: float, slideData: tuple, folder: str, size: tuple = SIZE, slideSize: tuple = (14.5, 7.5)):
    slideSizeX, slideSizeY = slideSize
    i, howMuchi = slideData
    imgInchSize = 4.5
    space, typeOfSpace, imgInchSize, heightOfImg = calculateSpace(
        slideSizeX, imgInchSize, howMuchi, size)
    graph = Image.open(graph_image_path)
    graph = graph.resize(size)
    graph_path = f"""{presentationsFolder}/{
        folder}/supportImage/resized_graph_{i}_{howMuchi}.png"""
    graph.save(graph_path)
    positionX = None
    if typeOfSpace == 'equi':
        positionX = space+((i-1)*(imgInchSize+space))
    else:
        positionX = space+((i-1)*imgInchSize)

    slide.shapes.add_picture(graph_path, Inches(positionX), Inches(Yposition),
                             width=Inches(imgInchSize), height=Inches(heightOfImg))
    if howMuchi == 2:
        space = 0.5
    doGreatComment(graph_note_path, (positionX, Yposition+heightOfImg+0.15,
                                     imgInchSize-0.03, slideSizeY-space-Yposition-heightOfImg), slide=slide)


def getPriorityText(priority: str):
    if priority == "0":
        priorityText = 'New Rs'
    if priority == "1":
        priorityText = 'Priority 1'
    elif priority == "2":
        priorityText = 'Priority 2'
    elif priority == "3":
        priorityText = 'Priority 3'
    elif priority == "4":
        priorityText = 'Priority 4'
    else:
        priorityText = 'Priority 5'
    return priorityText


def basicSlide(prs, country, actionTitle, priorityText):
    WIDTH, HEIGHT = (14.5, 7.5)
    slide_layout = prs.slide_layouts[6]  # Using a blank layout
    slide = prs.slides.add_slide(slide_layout)
    country_image_path = f'resources/{country}.png'
    didi = f'resources/didi.png'
    slide.shapes.add_picture(country_image_path, Inches(
        0.5), Inches(0.1), width=Inches(0.75), height=Inches(0.75))
    slide.shapes.add_picture(didi, Inches(
        13), Inches(0.01), width=Inches(1.5), height=Inches(0.84375))
    add_text(slide, actionTitle, font_size=0.5,
             bold=True, color=(252, 76, 2), position=(TITLE_START[0], TITLE_START[1], WIDTH-TITLE_START[0]-2-0.5, 0.75),)
    add_text(slide, f'Action Title', font_size=0.5*0.7,
             bold=True, color=(0, 0, 0), position=(ACTIONTITLE_START[0], ACTIONTITLE_START[1], WIDTH-0.5, 0.5), vertical=True)
    add_text(slide, f'{priorityText} - ', font_size=0.25,
             bold=False, color=(0, 0, 0), position=(SUBTITLE_START[0], SUBTITLE_START[1], WIDTH-0.5, 0.4), vertical=True)
    return slide


def makePresentation(MAC: bool = True, prioridades: list = ['1', '2', '3', '4',], bigFolder=plotsFolder, imagesPerSlide: int = 3):

    if MAC:
        paises = ['CO', 'PE', 'CR',]
    else:
        paises = ['MX',]
    createFolder(bigFolder)
    prs = Presentation()
    for pais in paises:
        for priority in prioridades:
            graph_image_paths, graph_text_paths = getImagesAndNotes(
                bigFolder, pais, priority)
            width, height = (14.5, 7.5)
            prs.slide_width = Inches(width)
            prs.slide_height = Inches(height)
            createFolder(f'{bigFolder}/supportImage')
            output_dir = f'{presentationsFolder}'
            for i in range(0, len(graph_image_paths)):
                if i % imagesPerSlide == 0:
                    priorityText = getPriorityText(priority)
                    title = ''
                    if i // imagesPerSlide == 0:
                        title = 'PERFORMANCE'
                    elif i // imagesPerSlide == 1:
                        title = 'SERVICE'
                    elif i // imagesPerSlide == 2:
                        title = 'EXPOSURE AND BURN'
                    else:
                        title = 'NO TITLE'
                    slide = basicSlide(prs, pais, title, priorityText)
                addGraphAndText(slide, graph_image_paths[i], graph_text_paths[i],
                                IMG_TOP, (i % imagesPerSlide+1, min(imagesPerSlide, len(graph_image_paths)-imagesPerSlide*(i//imagesPerSlide))), bigFolder)
                if i % imagesPerSlide == 0:
                    space, a, imgInchSize, heightOfImg = calculateSpace(width, 4.5, min(
                        imagesPerSlide, len(graph_image_paths)-imagesPerSlide*(i//imagesPerSlide)))
                    add_text(slide, f'Comments', font_size=0.2,
                             bold=False, color=(255, 255, 255), position=(space, IMG_TOP+heightOfImg-0.15, imgInchSize-0.03, 0.3), vertical=True,
                             BG=(252, 76, 2))
        # Country overview

        graph_image_paths, graph_text_paths = getImagesAndNotes(
            bigFolder, pais, priority, country=True)
        for i in range(0, len(graph_image_paths)):
            if i % 2 == 0:
                slide = basicSlide(prs, pais, 'PERFORMANCE',
                                   'Country overview')
                space, a, imgInchSize, heightOfImg = calculateSpace(
                    width, 4.5, 2)
                add_text(slide, f'Comments', font_size=0.2,
                         bold=False, color=(255, 255, 255), position=(space, IMG_TOP+heightOfImg-0.15, imgInchSize-0.03, 0.3), vertical=True,
                         BG=(252, 76, 2))
            addGraphAndText(slide, graph_image_paths[i], graph_text_paths[i],
                            IMG_TOP, ((i % 2)+1, 2), bigFolder, size=SIZE)

        print(f'Acabé con {pais} en {bigFolder}!')
    if MAC:
        add = 'MAC'
    else:
        add = 'MX'
    prs.save(f'{output_dir}/{add}_{bigFolder}.pptx')


def main():
    bigFolders = [plotsFolder,]
    for bigFolder in bigFolders:
        if MX:
            makePresentation(MAC=False, bigFolder=bigFolder)
        if MAC:
            makePresentation(MAC=True, bigFolder=bigFolder)


if __name__ == '__main__':
    main()
